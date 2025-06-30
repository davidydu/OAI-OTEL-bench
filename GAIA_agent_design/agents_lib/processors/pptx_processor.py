from __future__ import annotations

from pptx import Presentation

from ..file_router import Attachment


class PPTXProcessorAgent:
    """Extract text from PowerPoint presentations."""

    def process(self, att: Attachment) -> str:
        pres = Presentation(att.path)
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        texts.append(txt)
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = "\t".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            texts.append(row_text)
            if getattr(slide, "has_notes_slide", False):
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    texts.append(notes.strip())
        return "\n".join(texts)
